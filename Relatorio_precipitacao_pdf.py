1# -*- coding: utf-8 -*-
"""
Created on Wed Jan 24 14:11:45 2024

@author: FelipeSousa
"""
from selenium import webdriver
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.chrome.service import Service as ChromeService
from selenium.webdriver.common.by import By
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.common.keys import Keys
from pptx import Presentation
from datetime import date, timedelta
import os
from PIL import Image
import tempfile
import time
import comtypes.client

def cria_pasta(caminho):
    diretorio_base = caminho + r'Estudos/Previsão de Precipitação/Imagens'
    nome_pasta = date.today().strftime("%Y.%m.%d")
    caminho_pasta = os.path.join(diretorio_base, nome_pasta)
    caminho_pastas = []
    
    # Verifica se a pasta já existe antes de criar
    if not os.path.exists(caminho_pasta):
        os.makedirs(caminho_pasta)
    else:
        pass
    
    # Criação das subpastas dentro da pasta do dia
    subpastas = ["Precipitação Observada", "Previsão de Precipitação"]
    for subpasta in subpastas:
        caminho_subpasta = os.path.join(caminho_pasta, subpasta)
        
        if subpasta == r'Precipitação Observada':
            caminho_pastas.append(caminho_subpasta)
        else:
            pass
        
        # Verifica se a subpasta já existe antes de criar
        if not os.path.exists(caminho_subpasta):
            os.makedirs(caminho_subpasta)
        else:
            pass
        
    # Criação das subpastas dentro da pasta de Previsão de Precipitação
    subpastas_prev_prec = ["ECMWF", "ETA", "GEFS"]
    for subpasta in subpastas_prev_prec:
        caminho_subpasta_prev_prec = os.path.join(caminho_subpasta, subpasta)
        caminho_pastas.append(caminho_subpasta_prev_prec)
        
        # Verifica se a subpasta já existe antes de criar
        if not os.path.exists(caminho_subpasta_prev_prec):
            os.makedirs(caminho_subpasta_prev_prec)
        else:
            pass
    
    return caminho_pastas

def url_imagens():
    image_urls = []
    image_urls.append('https://sintegre.ons.org.br/sites/9/38/Documents/images/operacao_integrada/meteorologia/oshad_50_d.gif')
    
    for x in range(1,11):
        image_urls.append('https://sintegre.ons.org.br/sites/9/38/Documents/images/operacao_integrada/meteorologia/eta/shad50_'+str(x)+'.gif')
    
    for x in range(1,11):
        image_urls.append('https://sintegre.ons.org.br/sites/9/38/Documents/images/operacao_integrada/meteorologia/global/global50_50_'+str(x)+'.gif')
        
    for x in range(1,11):
        image_urls.append('https://sintegre.ons.org.br/sites/9/38/Documents/images/operacao_integrada/meteorologia/ECMWF/ecmwf50_'+str(x)+'.gif')
    
    return image_urls

def baixa_imagens(caminho_prec_obs, caminho_prev_ECMWF, caminho_prev_ETA, caminho_prev_GEFS, image_urls):
    options = webdriver.ChromeOptions()
    options.add_argument("--headless=new")
    servico = ChromeService(ChromeDriverManager().install())
    navegador = webdriver.Chrome(service=servico, options=options)
    wait = WebDriverWait(navegador, 20) 
    
    # Passo 1:
    navegador.get("https://sintegre.ons.org.br/")
    
    # Passo 2:
    element = wait.until(EC.element_to_be_clickable((By.XPATH, '//*[@id="form.username"]/input[1]')))
    navegador.find_element('xpath', '//*[@id="username"]').send_keys("PEDRO.RODRIGUES@NEOGIER.COM.BR")
    navegador.find_element('xpath', '//*[@id="form.username"]/input[1]').click()
    time.sleep(5)
    
    # Passo 3:
    element = wait.until(EC.element_to_be_clickable((By.XPATH, '//*[@id="form.password"]/input[1]')))
    navegador.find_element('xpath', '//*[@id="password"]').send_keys("Pc2599emp%")
    #navegador.find_element('xpath', '//*[@id="form.password"]/input[1]').click()
    navegador.find_element('xpath', '//*[@id="form.password"]/input[1]').send_keys(Keys.ENTER)
    time.sleep(5)
    
    # Passo 4:
    for i, image_url in enumerate(image_urls):
        navegador.get(image_url)
        
        # Encontrar a posição e tamanho do elemento que você deseja capturar
        element = navegador.find_element('xpath', '/html/body/img')
        location = element.location
        size = element.size
        
        # Captura uma screenshot da zona específica
        screenshot_path = 'screenshot_temp.png'
        navegador.save_screenshot(screenshot_path)

        # Abre a screenshot usando o Pillow (Python Imaging Library)
        image = Image.open(screenshot_path)
        
        # Crop na imagem para obter apenas a zona desejada
        left = location['x']
        top = location['y']
        right = location['x'] + size['width']
        bottom = location['y'] + size['height']
        image_cropped = image.crop((left, top, right, bottom))
        
        if i == 0:
            image_cropped.save(os.path.join(caminho_prec_obs, str('D-1'+'.gif')))
        
        elif i <= 10:
            image_cropped.save(os.path.join(caminho_prev_ETA, str('D'+str(i-1)+'.gif')))
            
        elif i <= 20:
            image_cropped.save(os.path.join(caminho_prev_GEFS, str('D'+str(i-11)+'.gif')))
            
        elif i <= 30:
            image_cropped.save(os.path.join(caminho_prev_ECMWF, str('D'+str(i-21)+'.gif')))

def get_left_position(shape):
    return shape.left

def get_top_position(shape):
    return shape.top

def ajusta_data_slide(data, caminho_slide_modelo):
    presentation = Presentation(caminho_slide_modelo)
    aux_8_10 = 0
    
    for n_slide in range(len(presentation.slides)):
        aux_1_7 = 0
        if n_slide < 7:
            slide = presentation.slides[n_slide]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if aux_1_7 == 0:
                                pass
                            elif aux_1_7 == 1:
                                run.text = (data + timedelta(days=aux_1_7-1)).strftime("%d/%m")
                            else:
                                run.text = (data + timedelta(days=aux_1_7-2)).strftime("%d/%m")
                        aux_1_7 += 1

        else:
            slide = presentation.slides[n_slide]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text == "PRECIPITAÇÃO (MM) ACUMULADA EM 24HS":
                                aux_8_10 -= 1
                            else:
                                run.text = (data + timedelta(days=aux_8_10-30)).strftime("%d/%m")
                    aux_8_10 += 1
    return (presentation)

def organiza_imagens(data, caminho_imagens):
    precipitacao_observada = {}
    previsao_precipitacao_GEFS = {}
    previsao_precipitacao_ETA = {}
    previsao_precipitacao_ECMWF = {}
    tam_img = (0, 59, 433, 503)
    aux = 0
    for delta in range(100):
        data_path = os.path.join(caminho_imagens, str((data - timedelta(days=delta)).strftime("%Y.%m.%d")))
        if os.path.exists(data_path):
            if delta < 10:
                #Para precipitacao observada
                img_p_obs_path = os.path.join(data_path, r'Precipitação Observada', r'D-1.gif')
                if os.path.exists(img_p_obs_path):
                    precipitacao_observada[str('D-'+str(delta))] = Image.open(img_p_obs_path).crop(tam_img)
                else:
                    precipitacao_observada[str('D-'+str(delta))] = 0
                
                #Para previsao precipitacao
                for x in range(aux,10):
                    img_p_GEFS_path = os.path.join(data_path, r'Previsão de Precipitação',r'GEFS', str('D'+str(x)+'.gif'))
                    img_p_ETA_path = os.path.join(data_path, r'Previsão de Precipitação',r'ETA', str('D'+str(x)+'.gif'))
                    img_p_ECMWF_path = os.path.join(data_path, r'Previsão de Precipitação',r'ECMWF', str('D'+str(x)+'.gif'))
                    if x == aux:
                        if os.path.exists(img_p_GEFS_path):
                            previsao_precipitacao_GEFS[str('D-'+str(delta))] = [Image.open(img_p_GEFS_path).crop(tam_img)]
                        else:
                            previsao_precipitacao_GEFS[str('D-'+str(delta))] = [0]
                            
                        if os.path.exists(img_p_ETA_path):
                            previsao_precipitacao_ETA[str('D-'+str(delta))] = [Image.open(img_p_ETA_path).crop(tam_img)]
                        else:
                            previsao_precipitacao_ETA[str('D-'+str(delta))] = [0]
                            
                        if os.path.exists(img_p_ECMWF_path):
                            previsao_precipitacao_ECMWF[str('D-'+str(delta))] = [Image.open(img_p_ECMWF_path).crop(tam_img)]
                        else:
                            previsao_precipitacao_ECMWF[str('D-'+str(delta))] = [0]
                    
                    else:
                        if os.path.exists(img_p_GEFS_path):
                            previsao_precipitacao_GEFS[str('D-'+str(delta))].append(Image.open(img_p_GEFS_path).crop(tam_img))
                        else:
                            previsao_precipitacao_GEFS[str('D-'+str(delta))].append(0)
                            
                        if os.path.exists(img_p_ETA_path):
                            previsao_precipitacao_ETA[str('D-'+str(delta))].append(Image.open(img_p_ETA_path).crop(tam_img))
                        else:
                            previsao_precipitacao_ETA[str('D-'+str(delta))].append(0)
                            
                        if os.path.exists(img_p_ECMWF_path):
                            previsao_precipitacao_ECMWF[str('D-'+str(delta))].append(Image.open(img_p_ECMWF_path).crop(tam_img))
                        else:
                            previsao_precipitacao_ECMWF[str('D-'+str(delta))].append(0)
                
            else:
                #Para precipitacao observada
                img_p_obs_path = os.path.join(data_path, r'Precipitação Observada', r'D-1.gif')
                if os.path.exists(img_p_obs_path):
                    precipitacao_observada[str('D-'+str(delta))] = Image.open(img_p_obs_path).crop(tam_img)
                else:
                    precipitacao_observada[str('D-'+str(delta))] = 0        
            
        else:
            if delta < 10:
                #Para precipitacao observada
                precipitacao_observada[str('D-'+str(delta))] = 0
            
                #zerar os 3 dicionários de previsao
                for x in range(aux,10):
                    if x == aux:
                        previsao_precipitacao_GEFS[str('D-'+str(delta))] = [0]
                        previsao_precipitacao_ETA[str('D-'+str(delta))] = [0]
                        previsao_precipitacao_ECMWF[str('D-'+str(delta))] = [0]
                    else:
                        previsao_precipitacao_GEFS[str('D-'+str(delta))].append(0)
                        previsao_precipitacao_ETA[str('D-'+str(delta))].append(0)
                        previsao_precipitacao_ECMWF[str('D-'+str(delta))].append(0)
            
            else:
                precipitacao_observada[str('D-'+str(delta))] = 0
        aux += 1
    
    return(precipitacao_observada, previsao_precipitacao_GEFS, previsao_precipitacao_ETA, previsao_precipitacao_ECMWF)

def ajusta_imagens_slide_1_6(data, presentation, precipitacao_observada, previsao_precipitacao_GEFS, previsao_precipitacao_ETA, previsao_precipitacao_ECMWF):
    for n_slide in range(0,6):      
        slide = presentation.slides[n_slide]
        shapes = []
        
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Verifique se o shape é uma imagem (código 13)
                shapes.append(shape)
        
        # Ordene as formas com base na posição vertical (de cima para baixo)
        shapes.sort(key = lambda shape: get_left_position(shape)**0.8 + get_top_position(shape))
        
        for i, shape in enumerate(shapes):
            largura = shape.width
            altura = shape.height
            
            if n_slide == 0:
                if i == 0:
                    if previsao_precipitacao_GEFS['D-9'][i] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-9'][i].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                elif i <= 2:
                    if previsao_precipitacao_GEFS['D-8'][i-1] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-8'][i-1].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 5:
                    if previsao_precipitacao_GEFS['D-7'][i-3] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-7'][i-3].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 9:
                    if previsao_precipitacao_GEFS['D-6'][i-6] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-6'][i-6].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                        
                elif i <= 14:
                    if previsao_precipitacao_GEFS['D-5'][i-10] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-5'][i-10].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                else:
                    pass  

            elif n_slide == 1:
                if i <= 5:
                    if previsao_precipitacao_GEFS['D-4'][i] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-4'][i].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                elif i <= 12:
                    if previsao_precipitacao_GEFS['D-3'][i-6] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-3'][i-6].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 20:
                    if previsao_precipitacao_GEFS['D-2'][i-13] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-2'][i-13].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 29:
                    if previsao_precipitacao_GEFS['D-1'][i-21] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-1'][i-21].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                        
                elif i <= 39:
                    if previsao_precipitacao_GEFS['D-0'][i-30] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_GEFS['D-0'][i-30].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                else:
                    pass           
            
            elif n_slide == 2:
                if i == 0:
                    if previsao_precipitacao_ETA['D-9'][0] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-9'][0].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                elif i <= 2:
                    if previsao_precipitacao_ETA['D-8'][i-1] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-8'][i-1].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 5:
                    if previsao_precipitacao_ETA['D-7'][i-3] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-7'][i-3].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 9:
                    if previsao_precipitacao_ETA['D-6'][i-6] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-6'][i-6].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                        
                elif i <= 14:
                    if previsao_precipitacao_ETA['D-5'][i-10] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-5'][i-10].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                else:
                    pass   

            elif n_slide == 3:
                if i <= 5:
                    if previsao_precipitacao_ETA['D-4'][i] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-4'][i].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                elif i <= 12:
                    if previsao_precipitacao_ETA['D-3'][i-6] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-3'][i-6].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 20:
                    if previsao_precipitacao_ETA['D-2'][i-13] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-2'][i-13].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 29:
                    if previsao_precipitacao_ETA['D-1'][i-21] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-1'][i-21].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                        
                elif i <= 39:
                    if previsao_precipitacao_ETA['D-0'][i-30] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ETA['D-0'][i-30].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                else:
                    pass      
                
            elif n_slide == 4:
                if i == 0:
                    if previsao_precipitacao_ECMWF['D-9'][0] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-9'][0].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                elif i <= 2:
                    if previsao_precipitacao_ECMWF['D-8'][i-1] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-8'][i-1].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 5:
                    if previsao_precipitacao_ECMWF['D-7'][i-3] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-7'][i-3].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 9:
                    if previsao_precipitacao_ECMWF['D-6'][i-6] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-6'][i-6].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                        
                elif i <= 14:
                    if previsao_precipitacao_ECMWF['D-5'][i-10] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-5'][i-10].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                else:
                    pass

            elif n_slide == 5:
                if i <= 5:
                    if previsao_precipitacao_ECMWF['D-4'][i] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-4'][i].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                elif i <= 12:
                    if previsao_precipitacao_ECMWF['D-3'][i-6] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-3'][i-6].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 20:
                    if previsao_precipitacao_ECMWF['D-2'][i-13] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-2'][i-13].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

                elif i <= 29:
                    if previsao_precipitacao_ECMWF['D-1'][i-21] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-1'][i-21].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                        
                elif i <= 39:
                    if previsao_precipitacao_ECMWF['D-0'][i-30] == 0:
                        pass
                    else:
                        # Insere a nova imagem nas mesmas coordenadas
                        shape.element.getparent().remove(shape.element)
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                            previsao_precipitacao_ECMWF['D-0'][i-30].save(temp_file.name)
                        slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
                
                else:
                    pass  

    return (presentation)

def ajusta_imagens_slide_7(data, presentation, precipitacao_observada, previsao_precipitacao_GEFS, previsao_precipitacao_ETA, previsao_precipitacao_ECMWF):
    slide = presentation.slides[6]
    shapes = []
    
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Verifique se o shape é uma imagem (código 13)
            shapes.append(shape)
    
    # Ordene as formas com base na posição vertical (de cima para baixo)
    shapes.sort(key = lambda shape: get_left_position(shape)**0.8 + get_top_position(shape))
    
    for i, shape in enumerate(shapes):
        largura = shape.width
        altura = shape.height
        
        if i < 10:
            if previsao_precipitacao_GEFS['D-0'][i] == 0:
                pass
            else:
                # Insere a nova imagem nas mesmas coordenadas
                shape.element.getparent().remove(shape.element)
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                    previsao_precipitacao_GEFS['D-0'][i].save(temp_file.name)
                slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
        elif i < 20:
            if previsao_precipitacao_ETA['D-0'][i-10] == 0:
                pass
            else:
                # Insere a nova imagem nas mesmas coordenadas
                shape.element.getparent().remove(shape.element)
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                    previsao_precipitacao_ETA['D-0'][i-10].save(temp_file.name)
                slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)

        elif i < 30:
            if previsao_precipitacao_ECMWF['D-0'][i-20] == 0:
                pass
            else:
                # Insere a nova imagem nas mesmas coordenadas
                shape.element.getparent().remove(shape.element)
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                    previsao_precipitacao_ECMWF['D-0'][i-20].save(temp_file.name)
                slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
        else:
            pass
        
    return (presentation)

def ajusta_imagens_slide_8_10(data, presentation, precipitacao_observada, previsao_precipitacao_GEFS, previsao_precipitacao_ETA, previsao_precipitacao_ECMWF):
    aux = 1
    for n_slide in range(7,10):      
        slide = presentation.slides[n_slide]
        shapes = []
        
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Verifique se o shape é uma imagem (código 13)
                shapes.append(shape)
        
        # Ordene as formas com base na posição vertical (de cima para baixo)
        shapes.sort(key = lambda shape: get_left_position(shape)**0.8 + get_top_position(shape))
        
        for i, shape in enumerate(shapes):
            largura = shape.width
            altura = shape.height
            
            if i < 10:
                if precipitacao_observada[str('D-'+str(30-aux))] == 0:
                    pass
                else:
                    # Insere a nova imagem nas mesmas coordenadas
                    shape.element.getparent().remove(shape.element)
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
                        precipitacao_observada[str('D-'+str(30-aux))].save(temp_file.name)
                    slide.shapes.add_picture(temp_file.name, shape.left, shape.top, width=largura, height=altura)
            else:
                aux -= 1
            aux += 1
    return (presentation)

def salva_apresentacao(presentation, data, caminho):
    pptx_path = caminho + r'Estudos/Previsão de Precipitação/Resultados/Relatório %s.pptx' % (data.strftime("%Y-%m-%d"))
    pdf_path = caminho + r'Estudos/Previsão de Precipitação/Resultados/Relatório %s.pdf' % (data.strftime("%Y-%m-%d"))
    
    # Salva como arquivo PowerPoint
    presentation.save(pptx_path)
    
    try:
        # Inicializa o PowerPoint via comtypes
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(pptx_path)
        
        # Salva como PDF
        presentation.SaveAs(pdf_path, 32)  # 32 é o formato PDF para PowerPoint
        print(f"Arquivo salvo com sucesso em: {pdf_path}")
        
    except Exception as e:
        print(f"Erro ao salvar o arquivo em PDF: {e}")
    
    finally:
        # Fecha o PowerPoint
        if 'presentation' in locals():
            presentation.Close()
        if 'powerpoint' in locals():
            powerpoint.Quit()

def iniciacao():
    print("Escolha uma opção:")
    print("1. C:/Users/FelipeSousa/OneDrive - NEOGIER ENERGIA/Neogier Y/")
    print("1. C:/Users/felip/OneDrive - NEOGIER ENERGIA/Neogier Y/")
    print("2. C:/Users/MARCOS/OneDrive - NEOGIER ENERGIA/Neogier Y/")
    print("3. C:/Users/mvgus/OneDrive - NEOGIER ENERGIA (1)/Neogier Y")
    opcao = int(input("Digite o número da opção: "))
    
    if opcao == 1:
        caminho = r"C:/Users/FelipeSousa/OneDrive - NEOGIER ENERGIA/Neogier Y/"
    
    if opcao == 2:
        caminho = r"C:/Users/felip/OneDrive - NEOGIER ENERGIA/Neogier Y/"
    
    if opcao == 3:
        caminho = r"C:/Users/MARCOS/OneDrive - NEOGIER ENERGIA/Neogier Y/"
    
    if opcao == 4:
        caminho = r"C:/Users/mvgus/OneDrive - NEOGIER ENERGIA (1)/Neogier Y/"
        
    return (caminho)

def main():
    caminho = iniciacao()
    caminho_prec_obs = cria_pasta(caminho)[0]
    caminho_prev_ECMWF = cria_pasta(caminho)[1]
    caminho_prev_ETA = cria_pasta(caminho)[2]
    caminho_prev_GEFS = cria_pasta(caminho)[3]
    image_urls = url_imagens()
    baixa_imagens(caminho_prec_obs, caminho_prev_ECMWF, caminho_prev_ETA, caminho_prev_GEFS, image_urls)
    data = date.today()
    caminho_slide_modelo = caminho + r'Estudos/Previsão de Precipitação/Programa/Slide Modelo.pptx'
    caminho_imagens = caminho + r'Estudos/Previsão de Precipitação/Imagens'
    precipitacao_observada = organiza_imagens(data, caminho_imagens)[0]
    previsao_precipitacao_GEFS = organiza_imagens(data, caminho_imagens)[1]
    previsao_precipitacao_ETA = organiza_imagens(data, caminho_imagens)[2]
    previsao_precipitacao_ECMWF = organiza_imagens(data, caminho_imagens)[3]
    apresentacao_data_ajustada = ajusta_data_slide(data, caminho_slide_modelo)
    apresentacao_slide_1_6_ajustada = ajusta_imagens_slide_1_6(data, apresentacao_data_ajustada, precipitacao_observada, previsao_precipitacao_GEFS, previsao_precipitacao_ETA, previsao_precipitacao_ECMWF)
    apresentacao_slide_7_ajustada = ajusta_imagens_slide_7(data, apresentacao_slide_1_6_ajustada, precipitacao_observada, previsao_precipitacao_GEFS, previsao_precipitacao_ETA, previsao_precipitacao_ECMWF)
    apresentacao_slide_8_10_ajustada = ajusta_imagens_slide_8_10(data, apresentacao_slide_7_ajustada, precipitacao_observada, previsao_precipitacao_GEFS, previsao_precipitacao_ETA, previsao_precipitacao_ECMWF)
    salva_apresentacao(apresentacao_slide_8_10_ajustada, data, caminho)
main()


